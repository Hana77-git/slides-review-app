import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path


def read_pptx(inputpath, outpath):
  prs = Presentation(inputpath)
  if not outpath.exists():
    outpath.mkdir(parents=True)
  
  for i, slide in enumerate(prs.slides):
    text_items = []
    image_count = 0

    for shape in slide.shapes:
        if not shape.has_text_frame:
            # 画像がある場合は保存
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_format = image.ext
                img_filename = f"slide{i+1}_img{image_count+1}.{image_format}"
                saveimgpath = outpath / img_filename
                with open(saveimgpath, "wb") as f:
                    f.write(image_bytes)
                image_count += 1
        else:
            # テキスト抽出
            text = shape.text_frame.text.strip()
            if text:
                text_items.append(text)

    # スライド単位でテキスト出力
    txtsavepath = outpath / f"slide{i+1}_text.txt"
    with open(txtsavepath, "w", encoding="utf-8") as f:
        f.write("\n\n".join(text_items))
